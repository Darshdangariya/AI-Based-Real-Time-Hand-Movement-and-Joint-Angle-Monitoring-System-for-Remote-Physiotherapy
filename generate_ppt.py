from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Cervical Posture Detection System"
    subtitle.text = "AI-Powered Real-time Posture Monitoring for Cervical Physiotherapy\n\nDeveloped by:\n23aiml011\n23aiml034\n23aiml035"

    # 2. Objectives
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title, body = slide.shapes.title, slide.placeholders[1]
    title.text = "Project Objectives"
    tf = body.text_frame
    tf.text = "Develop a real-time computer vision application for posture monitoring."
    p = tf.add_paragraph()
    p.text = "Provide immediate, clinical-grade feedback for cervical spine exercises."
    p = tf.add_paragraph()
    p.text = "Assist physiotherapy clinics and home rehabilitation programs."
    p = tf.add_paragraph()
    p.text = "Reduce the risk of cervical injuries due to poor posture."

    # 3. Methodology
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title, body = slide.shapes.title, slide.placeholders[1]
    title.text = "Methodology"
    tf = body.text_frame
    tf.text = "Computer Vision Pipeline:"
    p = tf.add_paragraph()
    p.text = "Pose Detection: Google's MediaPipe Pose for body landmarks."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Face Analysis: MediaPipe Face Mesh for detailed facial landmarks."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Real-time Processing: Streamlit WebRTC for live video streaming."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Exercise Detection Algorithms:"
    p = tf.add_paragraph()
    p.text = "Geometric calculations for tracking Range of Motion (ROM)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Research-based thresholds for assessing posture correctness."
    p.level = 1

    # 4. Work Completed
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title, body = slide.shapes.title, slide.placeholders[1]
    title.text = "Work Completed"
    tf = body.text_frame
    tf.text = "Implemented 5 Exercise Types with real-time feedback:"
    p = tf.add_paragraph()
    p.text = "Cervical Flexion, Extension, Lateral Tilt, Rotation, Chin Tuck."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Built a responsive and modern professional interface."
    p = tf.add_paragraph()
    p.text = "Integrated color-coded feedback (Excellent/Good/Poor)."
    p = tf.add_paragraph()
    p.text = "Developed real-time video processing capabilities."

    # 5. Expected Outcomes
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title, body = slide.shapes.title, slide.placeholders[1]
    title.text = "Expected Outcomes"
    tf = body.text_frame
    tf.text = "Physical Therapy: Enable guided, self-monitored exercise sessions."
    p = tf.add_paragraph()
    p.text = "Posture Correction: Assist individuals in maintaining healthy neck posture."
    p = tf.add_paragraph()
    p.text = "Rehabilitation: Provide accurate assessments for post-injury neck exercises."
    p = tf.add_paragraph()
    p.text = "Prevention: Preemptively address repetitive strain injuries."

    # 6. Team Contribution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title, body = slide.shapes.title, slide.placeholders[1]
    title.text = "Team Contribution"
    tf = body.text_frame
    tf.text = "23aiml011"
    p = tf.add_paragraph()
    p.text = "Core Computer Vision Pipeline & MediaPipe Integration"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "23aiml034"
    p = tf.add_paragraph()
    p.text = "Exercise Detection Algorithms & Range of Motion (ROM) Calculations"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "23aiml035"
    p = tf.add_paragraph()
    p.text = "Streamlit Web UI, WebRTC Integration & Cloud Deployment Strategies"
    p.level = 1

    # 7. Q&A / Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Thank You"
    subtitle.text = "Questions & Answers\n\nPrototype Demonstration"

    prs.save("Cervical_Posture_Detection_Presentation.pptx")
    print("Presentation generated successfully at Cervical_Posture_Detection_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
